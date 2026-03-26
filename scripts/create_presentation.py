"""
Generates the budget mobile app presentation as .pptx
Run: python scripts/create_presentation.py
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Colours matching the HTML (budget-presentation.html) ──────────────────────
BG_DEEP  = RGBColor(0x06, 0x09, 0x12)   # --bg-deep
BG_BASE  = RGBColor(0x0A, 0x0E, 0x1A)   # --bg-base
CARD_BG  = RGBColor(0x0D, 0x13, 0x1F)   # --bg-card
CARD_BG2 = RGBColor(0x11, 0x1A, 0x2E)   # --bg-card2
GOLD     = RGBColor(0xC9, 0xA8, 0x4C)   # --gold
GOLD_L   = RGBColor(0xE5, 0xC7, 0x6B)   # --gold-light
GOLD_D   = RGBColor(0x8B, 0x71, 0x28)   # --gold-dim
WHITE    = RGBColor(0xF0, 0xE8, 0xD5)   # --text-primary
GRAY     = RGBColor(0x8B, 0x80, 0x70)   # --text-secondary
DIM      = RGBColor(0x3A, 0x34, 0x28)   # --text-dim
GREEN    = RGBColor(0x4A, 0xDE, 0x80)
BLUE     = RGBColor(0x7D, 0xB4, 0xF0)
BORDER   = RGBColor(0x2A, 0x22, 0x10)   # --border-dim

SW, SH = 13.33, 7.5   # 16:9 widescreen inches

prs = Presentation()
prs.slide_width  = Inches(SW)
prs.slide_height = Inches(SH)
blank_layout = prs.slide_layouts[6]


# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(blank_layout)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG_DEEP
    return s


def box(slide, x, y, w, h, fill=None, border=None, border_w=Pt(1), radius=None):
    """Add a filled rectangle. Returns the shape."""
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = border_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, color=WHITE,
        align=PP_ALIGN.RIGHT, wrap=True, italic=False):
    """Add a text box."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = wrap
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def multiline(slide, lines, x, y, w, h, size=14, color=WHITE, align=PP_ALIGN.RIGHT, bold=False):
    """Add a text box with multiple lines (list of (text, bold, color, size))."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in lines:
        if isinstance(item, str):
            item = (item, bold, color, size)
        line_text, line_bold, line_color, line_size = item
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line_text
        run.font.size = Pt(line_size)
        run.font.bold = line_bold
        run.font.color.rgb = line_color
    return tb


def gold_line(slide, x, y, w):
    """Horizontal gold divider line."""
    line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    line.line.color.rgb = GOLD
    line.line.width = Pt(1.5)


def slide_label_and_title(slide, label, title, gold_part='', y_label=0.35, y_title=0.65):
    """Render the slide-label + h2 style header (right-aligned Hebrew)."""
    txt(slide, label, 0.4, y_label, 12.5, 0.4, size=11, color=GOLD_D, align=PP_ALIGN.RIGHT)
    if gold_part:
        # Render title with gold ending using a two-run paragraph
        tb = slide.shapes.add_textbox(Inches(0.4), Inches(y_title), Inches(12.5), Inches(0.6))
        tb.word_wrap = False
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r1 = p.add_run()
        r1.text = title + ' '
        r1.font.size = Pt(26)
        r1.font.bold = True
        r1.font.color.rgb = WHITE
        r2 = p.add_run()
        r2.text = gold_part
        r2.font.size = Pt(26)
        r2.font.bold = True
        r2.font.color.rgb = GOLD
    else:
        txt(slide, title, 0.4, y_title, 12.5, 0.6, size=26, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)
    gold_line(slide, 0.4, y_title + 0.65, 12.5)


def card(slide, x, y, w, h, icon, title, body, title_color=GOLD_L, icon_size=20):
    """Feature/security card with icon + title + body."""
    box(slide, x, y, w, h, fill=CARD_BG, border=BORDER)
    txt(slide, icon, x + 0.12, y + 0.12, 0.5, 0.4, size=icon_size, align=PP_ALIGN.LEFT)
    txt(slide, title, x + 0.12, y + 0.55, w - 0.24, 0.35, size=13, bold=True, color=title_color, align=PP_ALIGN.RIGHT)
    txt(slide, body,  x + 0.12, y + 0.88, w - 0.24, h - 1.0, size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()

# Gold corner decorations (top-left, bottom-right)
for cx, cy, cw, ch in [(0, 0, 1.5, 0.06), (0, 0, 0.06, 1.2),
                        (SW-1.5, SH-0.06, 1.5, 0.06), (SW-0.06, SH-1.2, 0.06, 1.2)]:
    b = box(s1, cx, cy, cw, ch, fill=GOLD_D)

# Emblem ⚡
txt(s1, '⚡', 6.0, 0.5, 1.3, 0.9, size=42, align=PP_ALIGN.CENTER)

# Slide label
txt(s1, 'מערכת ניהול תקציב', 0.5, 1.5, 12.3, 0.45, size=13, color=GOLD_D, align=PP_ALIGN.CENTER)

# Main title
txt(s1, 'מטבע הבזק', 0.5, 1.95, 12.3, 1.1, size=64, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Unit subtitle
txt(s1, 'אוגדה 99 · מחשוב ומידע', 0.5, 3.05, 12.3, 0.45, size=16, color=GRAY, align=PP_ALIGN.CENTER)

# Gold divider
gold_line(s1, 3.5, 3.6, 6.3)

# Tagline
txt(s1, 'מערכת ניהול תקציב דיגיטלית — בזמן אמת, מהניד, ללא נייר.',
    0.5, 3.75, 12.3, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)
txt(s1, 'מותאמת לצרכי אוגדה ומאפשרת שקיפות תקציבית מלאה לכל הדרגים.',
    0.5, 4.12, 12.3, 0.4, size=14, color=GRAY, align=PP_ALIGN.CENTER)

# Badges row
badge_data = [
    ('⚡ React Native', GOLD, RGBColor(0x1A, 0x14, 0x04)),
    ('🔷 Firebase Auth', BLUE, RGBColor(0x04, 0x0E, 0x1A)),
    ('📱 Android — זמין עכשיו', GREEN, RGBColor(0x04, 0x1A, 0x0A)),
    ('🔒 הצפנה E2E', GOLD_L, RGBColor(0x1A, 0x14, 0x04)),
]
bx = 1.7
for btxt, bc, bfill in badge_data:
    b = box(s1, bx, 4.75, 2.4, 0.45, fill=bfill, border=bc, border_w=Pt(1.2))
    txt(s1, btxt, bx + 0.1, 4.77, 2.2, 0.38, size=11, bold=True, color=bc, align=PP_ALIGN.CENTER)
    bx += 2.6


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — PROBLEM & SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
slide_label_and_title(s2, 'האתגר והפתרון', 'לפני ואחרי:', 'מהנייר למובייל')

# Before card
box(s2, 0.4, 1.5, 5.9, 4.3, fill=RGBColor(0x14, 0x06, 0x06), border=RGBColor(0x5A, 0x10, 0x10))
txt(s2, '📱', 0.55, 1.55, 0.6, 0.5, size=22, align=PP_ALIGN.LEFT)
txt(s2, 'לפני — כלים לא מתאימים', 1.2, 1.58, 4.9, 0.4, size=14, bold=True,
    color=RGBColor(0xF8, 0x71, 0x71), align=PP_ALIGN.RIGHT)
before_items = [
    'תקשורת תקציבית דרך וואטסאפ — ערוצים לא מאובטחים',
    'ניהול תקציב בגיליונות Excel מסורבלים',
    'קבצים מרובים, גרסאות שונות, אין אמת אחת',
    'מידע רגיש נשלח ללא הצפנה בקבוצות צ\'אט',
    'אין נראות בזמן אמת — מחכים לדוח הבא',
    'אישורים בעל פה, ללא תיעוד ומעקב',
]
for i, item in enumerate(before_items):
    txt(s2, f'✗  {item}', 0.55, 2.1 + i * 0.55, 5.6, 0.5, size=11,
        color=RGBColor(0xFC, 0xA5, 0xA5), align=PP_ALIGN.RIGHT)

# After card
box(s2, 6.6, 1.5, 6.33, 4.3, fill=RGBColor(0x04, 0x12, 0x0A), border=RGBColor(0x10, 0x52, 0x30))
txt(s2, '📱', 6.75, 1.55, 0.6, 0.5, size=22, align=PP_ALIGN.LEFT)
txt(s2, 'אחרי — מטבע הבזק', 7.4, 1.58, 5.4, 0.4, size=14, bold=True,
    color=RGBColor(0x4A, 0xDE, 0x80), align=PP_ALIGN.RIGHT)
after_items = [
    'בקשות דיגיטליות מהניד — בכל מקום ובכל זמן',
    'תמונת תקציב חיה לכל תפקיד בזמן אמת',
    'מסלול אישור אוטומטי עם התראות Push',
    'דוחות מיידיים, גרפים ו-KPI cards',
    'היסטוריה מלאה ומבוקרת בענן',
    'אישורים מרחוק — 24/7',
]
for i, item in enumerate(after_items):
    txt(s2, f'✓  {item}', 6.75, 2.1 + i * 0.55, 6.0, 0.5, size=11,
        color=RGBColor(0x86, 0xEF, 0xAC), align=PP_ALIGN.RIGHT)

# KPI strip
kpis = [('0', 'נייר — בקשות דיגיטליות'), ('12', 'תפקידים עם הרשאות'), ('24/7', 'גישה מכל מקום'), ('<5′', 'זמן הגשת בקשה')]
kx = 0.4
for val, lbl in kpis:
    box(s2, kx, 5.95, 2.98, 1.1, fill=CARD_BG, border=BORDER)
    txt(s2, val, kx, 6.0, 2.98, 0.6, size=26, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    txt(s2, lbl, kx, 6.57, 2.98, 0.4, size=9, color=GRAY, align=PP_ALIGN.CENTER)
    kx += 3.13


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — FEATURES (9 modules)
# ══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
slide_label_and_title(s3, 'יכולות המערכת', '9 מודולים מרכזיים —', 'הכל באפליקציה אחת')

features = [
    ('📊', 'ניהול תקציב בזמן אמת', 'מעקב תקציב חי לכל יחידה, חריגות עם התראות מיידיות'),
    ('🛒', 'בקשות רכש דיגיטליות', 'תהליך אישור מדורג ללא נייר: סמ"פ ← קס"א ← קצ"ב'),
    ('👤', 'תפקידים והרשאות', '12 תפקידים ייעודיים: מח"ט, סמח"ט, קצ"ב, ק"מ ועוד'),
    ('🔔', 'התראות חכמות', 'Push Notifications מותאמות לתפקיד וממתינות לפעולה'),
    ('📈', 'נתונים ויזואליים', 'גרפים אינטראקטיביים, KPI cards, מגמות בזמן'),
    ('📅', 'תוכנית עבודה שנתית', 'ניהול, מעקב ועדכון תוכנית עבודה מהניד'),
    ('🏆', 'דוחות רווחים ושיאים', 'מעקב חסכונות, הישגים תקציביים, עמידה ביעדים'),
    ('🤖', 'AI Insights', 'בינה מלאכותית לניתוח נתוני תקציב והמלצות חכמות'),
    ('🔒', 'אבטחה מתקדמת', 'Firebase Auth + E2E Encryption + Forward Secrecy'),
]

cols = 3
card_w = 4.0
card_h = 1.5
gap_x = 0.44
gap_y = 0.22
start_x = 0.4
start_y = 1.55

for i, (icon, title, body) in enumerate(features):
    col = i % cols
    row = i // cols
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)
    box(s3, cx, cy, card_w, card_h, fill=CARD_BG, border=BORDER)
    num_color = GOLD_D
    txt(s3, f'0{i+1}', cx + card_w - 0.55, cy + 0.08, 0.45, 0.3, size=9, color=num_color, align=PP_ALIGN.LEFT)
    txt(s3, icon, cx + 0.1, cy + 0.1, 0.55, 0.5, size=18, align=PP_ALIGN.LEFT)
    txt(s3, title, cx + 0.1, cy + 0.62, card_w - 0.2, 0.35, size=12, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
    txt(s3, body,  cx + 0.1, cy + 0.97, card_w - 0.2, 0.45, size=9.5, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — APPROVAL FLOW
# ══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
slide_label_and_title(s4, 'תהליך עבודה', 'בקשת שימוש בתקציב —', 'פשוט, מהיר, דיגיטלי')

# Three flow steps
steps = [
    ('👨‍💼', '1', 'המשתמש', 'מגיש בקשה לשימוש בכסף התקציבי — פרטים, סכום וסיבה', False),
    ('⭐',  '2', 'מאשר',   'קצין תקציבים / מאו"ג / סמאו"ג', True),
    ('✅',  '3', 'עדכון מיידי', 'המשתמש מקבל התראה בזמן אמת — אושר או נדחה', False),
]

step_w = 3.4
step_h = 2.8
step_y = 1.6
gap = 0.5
start = 0.5

for i, (icon, num, role, action, highlight) in enumerate(steps):
    sx = start + i * (step_w + gap + 0.8)
    if highlight:
        box(s4, sx, step_y, step_w, step_h, fill=RGBColor(0x14, 0x11, 0x04), border=GOLD_D)
    else:
        box(s4, sx, step_y, step_w, step_h, fill=CARD_BG, border=BORDER)
    txt(s4, icon, sx + step_w/2 - 0.3, step_y + 0.18, 0.6, 0.55, size=26, align=PP_ALIGN.CENTER)
    # Circle number
    nc = GOLD if highlight else GOLD_D
    txt(s4, num, sx + step_w/2 - 0.25, step_y + 0.82, 0.5, 0.38, size=15, bold=True, color=nc, align=PP_ALIGN.CENTER)
    rc = GOLD_L if highlight else WHITE
    txt(s4, role, sx + 0.1, step_y + 1.28, step_w - 0.2, 0.38, size=13, bold=True, color=rc, align=PP_ALIGN.CENTER)
    txt(s4, action, sx + 0.1, step_y + 1.7, step_w - 0.2, 0.9, size=10.5, color=GRAY, align=PP_ALIGN.CENTER)
    # Arrow between steps
    if i < 2:
        ax = sx + step_w + 0.1
        txt(s4, '←', ax, step_y + step_h/2 - 0.2, 0.6, 0.4, size=22, color=GOLD, align=PP_ALIGN.CENTER)

# Feature pills
feats = [('⚡ מהיר', 'ללא טפסים, ללא נייר'), ('📋 מבוקר', 'תיעוד עם חותמת זמן'), ('🔔 חכם', 'התראה אוטומטית')]
fx = 0.5
for ftitle, fbody in feats:
    box(s4, fx, 4.65, 4.1, 1.1, fill=RGBColor(0x0C, 0x0E, 0x14), border=BORDER)
    txt(s4, ftitle, fx + 0.1, 4.7, 3.9, 0.38, size=12, bold=True, color=GOLD_L, align=PP_ALIGN.CENTER)
    txt(s4, fbody,  fx + 0.1, 5.07, 3.9, 0.55, size=10, color=GRAY, align=PP_ALIGN.CENTER)
    fx += 4.4


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — ROLES
# ══════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
slide_label_and_title(s5, 'הרשאות ותפקידים', 'מבנה הרשאות —', 'כל תפקיד רואה רק את שלו')

# Hierarchy on left
hier = [
    ('מח"ט / סמח"ט', GOLD_L, GOLD_D, 0),
    ('מאו"ג / סמאו"ג', BLUE, RGBColor(0x05, 0x10, 0x20), 0.4),
    ('קצין תקציבים', RGBColor(0x5E, 0xEA, 0xD4), RGBColor(0x05, 0x18, 0x16), 0.8),
    ('קס"א', GRAY, CARD_BG, 1.0),
    ('סמ"פ / חייל', GRAY, CARD_BG, 1.0),
]
hy = 1.55
for name, tc, fc, indent in hier:
    bx_x = 0.4 + indent
    box(s5, bx_x, hy, 3.5 - indent*0.5, 0.48, fill=fc, border=tc, border_w=Pt(1))
    txt(s5, name, bx_x + 0.12, hy + 0.05, 3.3, 0.38, size=12, bold=True, color=tc, align=PP_ALIGN.CENTER)
    hy += 0.62

# Role chips on right (2 cols × 3 rows)
roles = [
    ('מח"ט', 'גישה מלאה לכל הנתונים', 'מלא'),
    ('מאו"ג', 'אישור בקשות + צפייה', 'גבוה'),
    ('קצין תקציבים', 'עריכת תקציב + אישורים', 'גבוה'),
    ('קס"א', 'צפייה ביחידה + הגשה', 'בינוני'),
    ('ק"מ', 'ניהול תקציב ציוד', 'בינוני'),
    ('סמ"פ', 'הגשת בקשות רכש', 'בסיסי'),
]
rx = 4.3
ry = 1.55
col = 0
for rname, rperm, rlevel in roles:
    cx = rx + col * 4.55
    box(s5, cx, ry, 4.3, 1.45, fill=CARD_BG, border=BORDER)
    txt(s5, rname, cx + 0.12, ry + 0.1, 4.0, 0.4, size=13, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
    txt(s5, rperm, cx + 0.12, ry + 0.5, 4.0, 0.4, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
    lcolor = GREEN if rlevel == 'מלא' else (BLUE if rlevel == 'גבוה' else (GOLD_D if rlevel == 'בינוני' else GRAY))
    txt(s5, rlevel, cx + 0.12, ry + 0.92, 4.0, 0.38, size=9.5, bold=True, color=lcolor, align=PP_ALIGN.RIGHT)
    col = 1 - col
    if col == 0:
        ry += 1.6


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — SECURITY (Enhanced with E2EE + Forward Secrecy)
# ══════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
slide_label_and_title(s6, 'אבטחת מידע', 'אבטחה ברמת WhatsApp —', 'Forward Secrecy מלא')

sec_items = [
    ('🔑', 'Firebase Authentication', 'כניסה עם JWT tokens + Custom Claims לתפקיד. Sessions מוצפנים ב-Android Keystore. איפוס אוטומטי בחצות.'),
    ('🔥', 'Firestore Security Rules', 'כל גישה למסד הנתונים מוגנת בחוקים מחמירים. משתמש רואה רק נתונים שלו — אין גישה צולבת בין יחידות.'),
    ('🔐', 'E2E Encryption — ECDH P-256', 'כל הודעת צ\'אט מוצפנת עם ECDH + HKDF-SHA-256 + AES-256-GCM. המפתח נגזר לכל הודעה בנפרד.'),
    ('⚡', 'Forward Secrecy — מפתח חד-פעמי', 'כל הודעה מקבלת Ephemeral Key pair חדש שנמחק מיד. גם אם המפתח הארוך-טווח נגנב — ההודעות הישנות מוגנות.'),
    ('🛡️', 'Root Detection', '5 בדיקות Native: build-tags, su binaries, אפליקציות root, system props, /system writable. גישה חסומה על מכשירים פרוצים.'),
    ('🔒', 'Network Security Config', 'רק אישורי CA ממערכת מותרים. ללא User CAs, ללא cleartext traffic. HTTPS בלבד. allowBackup=false.'),
]

card_w2 = 6.1
card_h2 = 1.55
for i, (icon, title, body) in enumerate(sec_items):
    col = i % 2
    row = i // 2
    cx = 0.4 + col * (card_w2 + 0.65)
    cy = 1.55 + row * (card_h2 + 0.18)
    box(s6, cx, cy, card_w2, card_h2, fill=CARD_BG, border=BORDER)
    txt(s6, icon, cx + 0.15, cy + 0.12, 0.55, 0.55, size=22, align=PP_ALIGN.LEFT)
    txt(s6, title, cx + 0.82, cy + 0.1, card_w2 - 0.97, 0.42, size=12.5, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
    txt(s6, body,  cx + 0.15, cy + 0.6, card_w2 - 0.3, 0.85, size=9.5, color=GRAY, align=PP_ALIGN.RIGHT)

# Banner
box(s6, 0.4, 6.45, 12.53, 0.72, fill=RGBColor(0x0A, 0x0C, 0x08), border=GOLD_D)
txt(s6, '🛡️', 0.55, 6.48, 0.55, 0.55, size=22, align=PP_ALIGN.LEFT)
txt(s6, 'ProGuard + R8 Obfuscation · Least Privilege · Android Keystore · הצפנה at-rest ו-in-transit',
    1.15, 6.52, 11.6, 0.55, size=11, bold=False, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — E2EE DEEP DIVE (NEW security slide)
# ══════════════════════════════════════════════════════════════════════════════
s7 = new_slide()
slide_label_and_title(s7, 'הצפנה קצה-לקצה', 'פרוטוקול ההצפנה —', 'Signal Protocol Level')

# Protocol steps
steps7 = [
    ('1', 'יצירת מפתח חד-פעמי', 'Alice מייצרת Ephemeral EC P-256 Key Pair טרי לכל הודעה'),
    ('2', 'ECDH Key Exchange', 'shared_secret = ECDH(eph_priv, Bob_long_term_pub)'),
    ('3', 'HKDF Key Derivation', 'aes_key = HKDF-SHA-256(shared_secret, info="budget-mobile-v2")'),
    ('4', 'AES-256-GCM Encrypt', 'ciphertext = AES-256-GCM(aes_key, random_iv, plaintext)'),
    ('5', 'Wire Format Storage', '[0x02 | ephPub(65B) | wrappedKey(40B) | iv(12B) | ciphertext]'),
    ('6', 'מחיקת המפתח', 'ephA_priv נמחק — Forward Secrecy בתוקף מכאן ואילך'),
]

sw7 = 3.85
sh7 = 1.38
gx = 0.62
gy = 0.22
sx7 = 0.4
sy7 = 1.55

for i, (num, title, desc) in enumerate(steps7):
    col = i % 3
    row = i // 2
    if i < 3:
        col = i
        row = 0
    else:
        col = i - 3
        row = 1
    cx = sx7 + col * (sw7 + gx)
    cy = sy7 + row * (sh7 + gy)
    is_key = (i == 5)
    fc = RGBColor(0x14, 0x11, 0x04) if is_key else CARD_BG
    bc = GOLD if is_key else BORDER
    box(s7, cx, cy, sw7, sh7, fill=fc, border=bc)
    nc = GOLD if is_key else GOLD_D
    txt(s7, num, cx + 0.15, cy + 0.1, 0.4, 0.38, size=13, bold=True, color=nc, align=PP_ALIGN.LEFT)
    txt(s7, title, cx + 0.1, cy + 0.45, sw7 - 0.2, 0.38, size=11.5, bold=True, color=GOLD_L if is_key else WHITE, align=PP_ALIGN.RIGHT)
    txt(s7, desc,  cx + 0.1, cy + 0.85, sw7 - 0.2, 0.45, size=9.5, color=GRAY, align=PP_ALIGN.RIGHT)

# Comparison vs WhatsApp
box(s7, 0.4, 4.62, 5.9, 1.65, fill=RGBColor(0x04, 0x12, 0x0A), border=RGBColor(0x10, 0x52, 0x30))
txt(s7, '✅  WhatsApp — Signal Protocol', 0.55, 4.68, 5.6, 0.42, size=12, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
txt(s7, 'Ephemeral ECDH · AES-256-GCM · Forward Secrecy · Double Ratchet', 0.55, 5.1, 5.6, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
txt(s7, '✅  מטבע הבזק — אותו פרוטוקול', 0.55, 5.48, 5.6, 0.42, size=12, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s7, 'Ephemeral ECDH P-256 · AES-256-GCM · Forward Secrecy · AES-KW Sender Copy', 0.55, 5.88, 5.6, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)

# Key storage box
box(s7, 6.55, 4.62, 6.38, 1.65, fill=CARD_BG, border=BORDER)
txt(s7, '🔑  אחסון מפתחות', 6.7, 4.68, 6.1, 0.42, size=12, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s7, 'מפתח פרטי → Android Keystore (חומרה)', 6.7, 5.1, 6.1, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
txt(s7, 'מפתח ציבורי → Firestore (מוצפן TLS)', 6.7, 5.48, 6.1, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)
txt(s7, 'Sessions → expo-secure-store (Keystore)', 6.7, 5.88, 6.1, 0.38, size=10, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — DISTRIBUTION
# ══════════════════════════════════════════════════════════════════════════════
s8 = new_slide()
slide_label_and_title(s8, 'הפצה מאובטחת', 'הפצה פנים-ארגונית —', 'ללא חנות ציבורית')

# Android card
box(s8, 0.4, 1.55, 5.9, 3.5, fill=RGBColor(0x04, 0x12, 0x08), border=RGBColor(0x1A, 0x7A, 0x4A))
txt(s8, '🤖', 0.55, 1.65, 0.7, 0.65, size=32, align=PP_ALIGN.LEFT)
txt(s8, 'Android — קובץ APK', 1.35, 1.7, 4.8, 0.45, size=16, bold=True, color=GREEN, align=PP_ALIGN.RIGHT)
txt(s8, 'שולחים קישור ישיר להורדת קובץ ה-APK לכל חייל עם מכשיר Android.\nהתקנה פשוטה ישירות מהקישור — ללא Google Play, ללא חשיפה לציבור.\nהאפליקציה נשארת בתוך הארגון בלבד.',
    0.55, 2.22, 5.6, 1.6, size=11, color=GRAY, align=PP_ALIGN.RIGHT)
box(s8, 0.55, 3.95, 2.2, 0.55, fill=RGBColor(0x08, 0x28, 0x14), border=RGBColor(0x1A, 0x7A, 0x4A))
txt(s8, '● מוכן להפצה ✓', 0.65, 4.0, 2.0, 0.38, size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# iOS card
box(s8, 6.73, 1.55, 6.0, 3.5, fill=RGBColor(0x05, 0x08, 0x16), border=RGBColor(0x15, 0x35, 0x80))
txt(s8, '🍎', 6.88, 1.65, 0.7, 0.65, size=32, align=PP_ALIGN.LEFT)
txt(s8, 'iOS — Apple TestFlight', 7.68, 1.7, 4.9, 0.45, size=16, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)
txt(s8, 'שולחים הזמנה אישית לאייפון של כל חייל דרך Apple TestFlight.\nהמשתמש מותקן כ-Tester מורשה — רק מי שמוזמן יכול להוריד.\nאין גישה ציבורית, אין חשיפה ב-App Store.',
    6.88, 2.22, 5.7, 1.6, size=11, color=GRAY, align=PP_ALIGN.RIGHT)
box(s8, 6.88, 3.95, 2.5, 0.55, fill=RGBColor(0x05, 0x0C, 0x28), border=RGBColor(0x15, 0x35, 0x80))
txt(s8, '● עד 10,000 משתמשים', 7.0, 4.0, 2.3, 0.38, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Why not store
box(s8, 0.4, 5.35, 12.53, 1.3, fill=RGBColor(0x0C, 0x0E, 0x08), border=GOLD_D)
txt(s8, '🔒', 0.55, 5.42, 0.55, 0.55, size=22, align=PP_ALIGN.LEFT)
txt(s8, 'למה לא פרסום בחנות?', 1.2, 5.42, 11.5, 0.42, size=13, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s8, 'המערכת מכילה מידע ארגוני-ביטחוני — פרסום ב-App Store חושף לציבור הרחב. בשיטת TestFlight + APK ישיר, רק מי שאנחנו מאשרים מקבל גישה.',
    1.2, 5.85, 11.5, 0.65, size=10.5, color=GRAY, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — FUTURE
# ══════════════════════════════════════════════════════════════════════════════
s9 = new_slide()
slide_label_and_title(s9, 'הצעד הבא', 'העתיד —', 'אינטגרציה עם מחשבי צה"ל')

# IDF computers
box(s9, 0.4, 1.55, 5.9, 3.5, fill=CARD_BG, border=BORDER)
txt(s9, '🖥️', 0.55, 1.65, 0.8, 0.75, size=38, align=PP_ALIGN.LEFT)
txt(s9, 'שלב א׳ — מחשבי צה"ל', 1.45, 1.72, 4.7, 0.45, size=15, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s9, 'השלב הבא הוא להעלות את המערכת על מחשבי צה"ל הארגוניים — גישה ממחשב נייח בבסיס, בנוסף לאפליקציה הניידת. הממשק יותאם לשימוש דסקטופ.',
    0.55, 2.28, 5.6, 1.6, size=11, color=GRAY, align=PP_ALIGN.RIGHT)
box(s9, 0.55, 3.95, 1.8, 0.55, fill=RGBColor(0x05, 0x0C, 0x28), border=RGBColor(0x15, 0x35, 0x80))
txt(s9, 'בתכנון', 0.65, 4.0, 1.6, 0.38, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Azure AD
box(s9, 6.73, 1.55, 6.0, 3.5, fill=RGBColor(0x0C, 0x0E, 0x16), border=RGBColor(0x2D, 0x5B, 0xE3, ))
txt(s9, '🔷', 6.88, 1.65, 0.8, 0.75, size=38, align=PP_ALIGN.LEFT)
txt(s9, 'שלב ב׳ — Azure Active Directory', 7.78, 1.72, 4.8, 0.45, size=15, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s9, 'לאחר האינטגרציה עם מחשבי צה"ל, נחבר את המערכת ל-Azure AD של Microsoft — כניסה אוטומטית עם חשבון הצה"ל, ללא שם משתמש וסיסמה נפרדים. Single Sign-On עם הגנת MFA.',
    6.88, 2.28, 5.7, 1.6, size=11, color=GRAY, align=PP_ALIGN.RIGHT)
box(s9, 6.88, 3.95, 2.0, 0.55, fill=RGBColor(0x05, 0x0C, 0x28), border=RGBColor(0x15, 0x35, 0x80))
txt(s9, 'שלב עתידי', 7.0, 4.0, 1.8, 0.38, size=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Roadmap banner
box(s9, 0.4, 5.35, 12.53, 1.3, fill=RGBColor(0x0A, 0x0C, 0x08), border=GOLD_D)
txt(s9, '🛡️', 0.55, 5.42, 0.55, 0.55, size=22, align=PP_ALIGN.LEFT)
txt(s9, 'המסלול המלא:', 1.2, 5.42, 11.5, 0.42, size=13, bold=True, color=GOLD_L, align=PP_ALIGN.RIGHT)
txt(s9, 'כרגע — אפליקציה ניידת עם כניסה מאובטחת  ←  הרחבה למחשבי צה"ל  ←  כניסה אוטומטית עם חשבון Azure AD הארגוני',
    1.2, 5.85, 11.5, 0.65, size=11, color=GRAY, align=PP_ALIGN.RIGHT)

# Badges
badge_data9 = [
    ('📱 Android — מוכן עכשיו', GREEN, RGBColor(0x04, 0x12, 0x06)),
    ('🍎 iOS TestFlight — מוכן עכשיו', GREEN, RGBColor(0x04, 0x12, 0x06)),
    ('🖥️ מחשבי צה"ל — בתכנון', BLUE, RGBColor(0x04, 0x06, 0x14)),
    ('🔷 Azure AD — שלב עתידי', GOLD, RGBColor(0x12, 0x0E, 0x04)),
]
bx9 = 0.8
for btxt9, bc9, bfill9 in badge_data9:
    b9 = box(s9, bx9, 6.92, 2.9, 0.42, fill=bfill9, border=bc9, border_w=Pt(1))
    txt(s9, btxt9, bx9 + 0.08, 6.94, 2.74, 0.35, size=10, bold=True, color=bc9, align=PP_ALIGN.CENTER)
    bx9 += 3.1


# ══════════════════════════════════════════════════════════════════════════════
#  Save
# ══════════════════════════════════════════════════════════════════════════════
out_path = r'C:\Users\elads\OneDrive\Desktop\budget_presentation_v2.pptx'
prs.save(out_path)

msg = b'\nSaved: ' + out_path.encode('utf-8') + b'\n'
sys.stdout.buffer.write(msg)
